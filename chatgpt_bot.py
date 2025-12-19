import os
import io
from telegram import Update, InputFile
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, ContextTypes, filters

from openai import OpenAI
import google.generativeai as genai

from PyPDF2 import PdfReader
from pptx import Presentation
import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# ===== Keys =====
TELEGRAM_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
OPENAI_KEY = os.getenv("OPENAI_API_KEY")
GEMINI_KEY = os.getenv("GEMINI_API_KEY")

openai_client = OpenAI(api_key=OPENAI_KEY)
genai.configure(api_key=GEMINI_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-flash")

SYSTEM_PROMPT = (
    "You are CyberBuddy, a helpful cyber & tech assistant. "
    "Never say you are ChatGPT, GPT, Gemini, or an OpenAI/Google model. "
    "Always introduce yourself as CyberBuddy."
)

# 🧠 Memory per user
user_memory = {}
MAX_HISTORY = 12

# user -> selected model
user_model = {}

# ===== Helpers =====
def ask_gpt(history):
    resp = openai_client.chat.completions.create(
        model="gpt-4.1-mini",
        messages=[{"role": "system", "content": SYSTEM_PROMPT}] + history,
    )
    return resp.choices[0].message.content

def ask_gemini(history):
    text = SYSTEM_PROMPT + "\n"
    for m in history:
        text += f"{m['role'].upper()}: {m['content']}\n"
    resp = gemini_model.generate_content(text)
    return resp.text

def ai_reply(user_id):
    history = user_memory.get(user_id, [])
    model = user_model.get(user_id, "gpt")
    return ask_gpt(history) if model == "gpt" else ask_gemini(history)

# ===== Commands =====
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    user_model[uid] = "gpt"
    user_memory[uid] = []
    await update.message.reply_text(
        "👋 I’m CyberBuddy 🤖\n"
        "Ab main hamari baatein yaad rakhunga.\n\n"
        "🤖 AI switch:\n"
        "/use gpt\n"
        "/use gemini\n\n"
        "📝 Files bhejo ya chat karo.\n"
        "Memory clear: /reset"
    )

async def reset(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    user_memory[uid] = []
    await update.message.reply_text("🔄 Tumhari memory reset ho gayi!")

async def use_model(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    choice = context.args[0].lower() if context.args else ""
    if choice not in ["gpt", "gemini"]:
        await update.message.reply_text("Use: /use gpt  or  /use gemini")
        return
    user_model[uid] = choice
    await update.message.reply_text(f"✅ Now using {choice.upper()} with memory")

# ===== Text Handler with Memory =====
async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    text = update.message.text

    history = user_memory.get(uid, [])
    history.append({"role": "user", "content": text})
    history = history[-MAX_HISTORY:]

    try:
        reply = ai_reply(uid)
        await update.message.reply_text(reply)

        history.append({"role": "assistant", "content": reply})
        history = history[-MAX_HISTORY:]
        user_memory[uid] = history

    except Exception as e:
        print("Chat error:", e)
        await update.message.reply_text("⚠️ Kuch problem aa gayi.")

# ===== File Handler with Memory =====
async def handle_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    doc = update.message.document
    file = await context.bot.get_file(doc.file_id)
    path = f"temp_{doc.file_name}"
    await file.download_to_drive(path)

    extracted = ""
    try:
        if doc.file_name.lower().endswith(".pdf"):
            reader = PdfReader(path)
            for p in reader.pages:
                extracted += p.extract_text() or ""

        elif doc.file_name.lower().endswith(".pptx"):
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted += shape.text + "\n"

        elif doc.file_name.lower().endswith((".xlsx", ".xls")):
            df = pd.read_excel(path)
            extracted = df.head(20).to_string()

        else:
            await update.message.reply_text("❌ Unsupported file type.")
            return

        history = user_memory.get(uid, [])
        history.append({"role": "user", "content": f"File content:\n{extracted[:3000]}"})
        history = history[-MAX_HISTORY:]

        reply = ai_reply(uid)
        await update.message.reply_text(reply)

        history.append({"role": "assistant", "content": reply})
        history = history[-MAX_HISTORY:]
        user_memory[uid] = history

    finally:
        if os.path.exists(path):
            os.remove(path)

# ===== Create PDF =====
async def make_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    text = " ".join(context.args)
    if not text:
        await update.message.reply_text("Use: /makepdf your text")
        return

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.drawString(40, 750, text[:1000])
    c.save()
    buf.seek(0)

    await update.message.reply_document(InputFile(buf, filename="cyberbuddy.pdf"))

# ===== Main =====
def main():
    app = ApplicationBuilder().token(TELEGRAM_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("reset", reset))
    app.add_handler(CommandHandler("use", use_model))
    app.add_handler(CommandHandler("makepdf", make_pdf))

    app.add_handler(MessageHandler(filters.Document.ALL, handle_file))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

    print("🤖 CyberBuddy with memory + GPT + Gemini is running...")
    app.run_polling()

if __name__ == "__main__":
    main()
